"""
File Store — RAG over local files.
Indexes text from PC files (PDF, DOCX, TXT, code files) using Gemini
text-embedding-004 so Jarvis can search and retrieve relevant file content
in any conversation.

Storage:
  memory/file_embeddings/
    index.json   — [{id, file_path, file_name, file_hash, chunk_idx,
                      total_chunks, date, text, size_kb}]
    vectors.npy  — float32 array shape (N, 768)
"""

import hashlib
import json
import logging
import os
import sys
import threading
import time
from datetime import datetime
from pathlib import Path
from typing import Generator

import numpy as np

logger = logging.getLogger("jarvis.filestore")
_lock = threading.Lock()

EMBED_DIM   = 768
CHUNK_WORDS = 400   # words per chunk
CHUNK_OVER  = 50    # word overlap between adjacent chunks
MAX_FILE_MB = 15    # skip files larger than this

# ── Supported file types ──────────────────────────────────────────────────────

TEXT_SUFFIXES = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".html", ".css", ".json", ".xml", ".csv", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".sh", ".bat", ".ps1",
    ".c", ".cpp", ".h", ".java", ".go", ".rs", ".rb", ".php",
    ".sql", ".r", ".env", ".log", ".rst", ".tex",
}

BINARY_SUFFIXES = {".pdf", ".docx", ".doc", ".odt", ".rtf", ".pptx", ".xlsx"}

ALL_SUFFIXES = TEXT_SUFFIXES | BINARY_SUFFIXES

# Folders to skip when scanning
SKIP_DIRS = {
    "__pycache__", ".git", "node_modules", ".venv", "venv", "env",
    "dist", "build", ".next", ".cache", "site-packages",
    "AppData", "Windows", "System32", "Program Files",
    "Program Files (x86)", "$Recycle.Bin",
    # Sensitive — credentials, keys, browser passwords
    ".ssh", ".gnupg", ".aws", ".azure",
    "Credentials", "Protect", "Firefox", "Chrome", "Edge",
}

# Default folders indexed on first startup — broad coverage of user files
_HOME = Path.home()
DEFAULT_FOLDERS = [f for f in [
    _HOME / "Desktop",
    _HOME / "Documents",
    _HOME / "Downloads",
    _HOME / "OneDrive",
    _HOME / "OneDrive - Personal",
    _HOME / "Pictures",
    _HOME / "Videos",
    _HOME / "Music",
    _HOME / "Projects",
    _HOME / "Work",
    _HOME / "Code",
    _HOME / "Dev",
    _HOME / "repos",
] if f.exists()]


def _base_dir() -> Path:
    if getattr(sys, "frozen", False):
        return Path(sys.executable).parent
    return Path(__file__).resolve().parent.parent


_EMBED_DIR  = _base_dir() / "memory" / "file_embeddings"
_INDEX_PATH = _EMBED_DIR / "index.json"
_VEC_PATH   = _EMBED_DIR / "vectors.npy"


# ── I/O ───────────────────────────────────────────────────────────────────────

def _load_index() -> list:
    try:
        return json.loads(_INDEX_PATH.read_text(encoding="utf-8")) if _INDEX_PATH.exists() else []
    except Exception:
        return []


def _load_vectors() -> np.ndarray:
    try:
        return np.load(str(_VEC_PATH)).astype(np.float32) if _VEC_PATH.exists() else np.empty((0, EMBED_DIM), dtype=np.float32)
    except Exception:
        return np.empty((0, EMBED_DIM), dtype=np.float32)


def _save_index(index: list) -> None:
    _EMBED_DIR.mkdir(parents=True, exist_ok=True)
    _INDEX_PATH.write_text(json.dumps(index, indent=2, ensure_ascii=False), encoding="utf-8")


def _save_vectors(vecs: np.ndarray) -> None:
    _EMBED_DIR.mkdir(parents=True, exist_ok=True)
    tmp = _VEC_PATH.with_name("vectors.tmp.npy")
    np.save(str(tmp), vecs)
    os.replace(tmp, _VEC_PATH)


# ── Text extraction ───────────────────────────────────────────────────────────

def _file_hash(path: Path) -> str:
    """MD5 of first 128 KB — fast dedup without reading whole file."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        h.update(f.read(131072))
    return h.hexdigest()


def _extract_text(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix in TEXT_SUFFIXES:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""

    if suffix == ".pdf":
        # Try pdfplumber first, then PyPDF2
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages)
        except ImportError:
            pass
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            return ""

    if suffix in {".docx", ".doc", ".odt"}:
        try:
            import docx
            doc = docx.Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return ""

    if suffix == ".rtf":
        try:
            raw = path.read_text(encoding="utf-8", errors="ignore")
            import re
            raw = re.sub(r'\\[a-z]+\d* ?', ' ', raw)
            raw = re.sub(r'[{}\\]', ' ', raw)
            return raw
        except Exception:
            return ""

    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception:
            return ""

    if suffix == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        rows.append(line)
            return "\n".join(rows)
        except Exception:
            return ""

    return ""


def _chunk_text(text: str) -> list[str]:
    """Split text into overlapping word chunks."""
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    while start < len(words):
        end   = min(start + CHUNK_WORDS, len(words))
        chunk = " ".join(words[start:end])
        if len(chunk.strip()) >= 80:
            chunks.append(chunk)
        start += CHUNK_WORDS - CHUNK_OVER
    return chunks


# ── Embedding ─────────────────────────────────────────────────────────────────

def _embed(text: str, api_key: str) -> np.ndarray | None:
    try:
        from google import genai
        client   = genai.Client(api_key=api_key)
        response = client.models.embed_content(
            model="text-embedding-004",
            contents=text[:2000],
        )
        return np.array(response.embeddings[0].values, dtype=np.float32)
    except Exception as e:
        logger.error("Embed error: %s", e)
        return None


# ── Index a single file ───────────────────────────────────────────────────────

def index_file(path: Path, api_key: str, existing_hashes: set,
               existing_mtimes: dict | None = None) -> int:
    """
    Index a single file into the vector store.
    Skips if unchanged (same mtime+size, or same hash) or too large.
    Returns number of new chunks added.
    """
    try:
        if not path.is_file():
            return 0
        stat       = path.stat()
        size_bytes = stat.st_size
        mtime      = stat.st_mtime
        if size_bytes > MAX_FILE_MB * 1024 * 1024:
            return 0

        # Fast skip: same path, same mtime and size → definitely unchanged
        str_path = str(path)
        if existing_mtimes and str_path in existing_mtimes:
            cached = existing_mtimes[str_path]
            if cached.get("mtime") == mtime and cached.get("size") == size_bytes:
                return 0

        fhash = _file_hash(path)
        if fhash in existing_hashes:
            return 0  # same content hash → unchanged

        text = _extract_text(path)
        if not text or len(text.strip()) < 60:
            return 0

        chunks = _chunk_text(text)
        if not chunks:
            return 0

        date        = datetime.now().strftime("%Y-%m-%d")
        new_entries = []
        new_vecs    = []

        # Embed all chunks in batches of 20 (one API call per batch)
        BATCH = 20
        for batch_start in range(0, len(chunks), BATCH):
            batch_chunks = chunks[batch_start: batch_start + BATCH]
            try:
                from memory.semantic_store import get_embeddings_batch
                batch_vecs = get_embeddings_batch(batch_chunks, api_key)
            except Exception:
                batch_vecs = [_embed(c, api_key) for c in batch_chunks]

            for j, (chunk, vec) in enumerate(zip(batch_chunks, batch_vecs)):
                i = batch_start + j
                if vec is None:
                    continue
                new_entries.append({
                    "id":           f"file_{fhash[:8]}_{i}",
                    "file_path":    str(path),
                    "file_name":    path.name,
                    "file_ext":     path.suffix.lower(),
                    "file_hash":    fhash,
                    "chunk_idx":    i,
                    "total_chunks": len(chunks),
                    "date":         date,
                    "size_kb":      round(size_bytes / 1024, 1),
                    "mtime":        mtime,
                    "size":         size_bytes,
                    "text":         chunk[:500],
                })
                new_vecs.append(vec)
            time.sleep(0.3)   # pause between batches

        if not new_entries:
            return 0

        with _lock:
            index = _load_index()
            # Remove old chunks from same file (re-index case)
            index = [e for e in index if e.get("file_hash") != fhash]
            vecs  = _load_vectors()
            for entry, vec in zip(new_entries, new_vecs):
                index.append(entry)
                if vecs.shape[0] == 0:
                    vecs = vec.reshape(1, -1)
                else:
                    vecs = np.vstack([vecs, vec.reshape(1, -1)])
            _save_index(index)
            _save_vectors(vecs)

        logger.info("Indexed %d chunks ← %s", len(new_entries), path.name)
        return len(new_entries)

    except Exception as e:
        logger.error("index_file error (%s): %s", path.name, e)
        return 0


# ── Folder scanner ────────────────────────────────────────────────────────────

def _iter_files(folders: list[Path]) -> Generator[Path, None, None]:
    seen: set[Path] = set()
    for folder in folders:
        if not folder.exists():
            continue
        try:
            for path in folder.rglob("*"):
                if path in seen:
                    continue
                seen.add(path)
                if not path.is_file():
                    continue
                if any(skip in path.parts for skip in SKIP_DIRS):
                    continue
                if path.suffix.lower() in ALL_SUFFIXES:
                    yield path
        except PermissionError:
            continue


def scan_and_index(
    folders: list[str | Path],
    api_key: str,
    progress_callback=None,
) -> int:
    """
    Walk the given folders and index any new or changed files.
    Uses mtime+size cache to skip unchanged files without hashing.
    Rate-limited — safe to call in a background thread.

    progress_callback(files_done, total_files) — called every 50 files.
    Returns total chunks added.
    """
    folder_paths = [Path(f).expanduser().resolve() for f in folders]

    with _lock:
        existing = _load_index()

    existing_hashes = {e["file_hash"] for e in existing}
    # Build mtime cache: {file_path: {mtime, size}} for fast skip
    existing_mtimes: dict = {}
    for e in existing:
        fp = e.get("file_path", "")
        if fp and "mtime" in e and fp not in existing_mtimes:
            existing_mtimes[fp] = {"mtime": e["mtime"], "size": e["size"]}

    # Pre-count files so progress is meaningful
    all_files = list(_iter_files(folder_paths))
    total_files = len(all_files)

    total = 0
    for files_done, path in enumerate(all_files, start=1):
        added = index_file(path, api_key, existing_hashes, existing_mtimes)
        if added:
            total += added
            try:
                existing_hashes.add(_file_hash(path))
                st = path.stat()
                existing_mtimes[str(path)] = {"mtime": st.st_mtime, "size": st.st_size}
            except Exception:
                pass
        if progress_callback and files_done % 50 == 0:
            try:
                progress_callback(files_done, total_files)
            except Exception:
                pass

    logger.info("Scan done — %d new chunks across %d folder(s)", total, len(folder_paths))
    return total


def scan_default_folders(api_key: str, progress_callback=None) -> int:
    """Index Desktop, Documents, Downloads in the background."""
    return scan_and_index(DEFAULT_FOLDERS, api_key, progress_callback=progress_callback)


# ── Search ────────────────────────────────────────────────────────────────────

def _safe_query(query: str, max_len: int = 500) -> str:
    """Strip control characters and limit length of a search query."""
    import unicodedata
    cleaned = "".join(c for c in query if not unicodedata.category(c).startswith("C"))
    return cleaned[:max_len].strip()


_BLOCKED_PATH_SEGMENTS = {
    # Windows system dirs
    "windows", "system32", "syswow64", "program files", "program files (x86)",
    # Credential / key stores
    ".ssh", ".gnupg", ".aws", ".azure",
    "microsoft\\credentials", "microsoft\\protect",
    "roaming\\microsoft\\windows\\credentials",
    # Browser saved passwords
    "google\\chrome\\user data\\default",
    "microsoft\\edge\\user data\\default",
    "mozilla\\firefox",
    # Unix system dirs
    "etc", "usr", "bin", "sbin",
    # Jarvis own config (api keys)
    "config",
}


def _safe_folder_path(path_str: str) -> Path:
    """
    Resolve and validate a user-supplied folder path.
    Raises ValueError if the path contains any sensitive or system segment.
    """
    resolved = Path(path_str).expanduser().resolve()
    lower = str(resolved).lower().replace("/", "\\")
    for segment in _BLOCKED_PATH_SEGMENTS:
        if segment in lower:
            raise ValueError(f"Access denied — sensitive path blocked: {resolved}")
    return resolved


def search_files(query: str, api_key: str, top_k: int = 5) -> list[dict]:
    """
    Semantic search over indexed files.
    Returns one best chunk per file, sorted by cosine similarity.
    """
    query = _safe_query(query)
    if not query:
        return []
    try:
        with _lock:
            index = _load_index()
            vecs  = _load_vectors()

        if not index or vecs.shape[0] == 0:
            return []

        q_vec = _embed(query, api_key)
        if q_vec is None:
            return []

        norms  = np.linalg.norm(vecs, axis=1)
        norms[norms == 0] = 1e-9
        q_norm = np.linalg.norm(q_vec) or 1e-9
        scores = (vecs @ q_vec) / (norms * q_norm)

        top_indices = np.argsort(scores)[::-1]
        results: list[dict] = []
        seen_files: set[str] = set()

        for i in top_indices:
            if i >= len(index):
                continue
            entry = dict(index[i])
            entry["score"] = float(scores[i])
            fpath = entry.get("file_path", "")
            if fpath not in seen_files:
                seen_files.add(fpath)
                results.append(entry)
            if len(results) >= top_k:
                break

        return results
    except Exception as e:
        logger.error("search error: %s", e)
        return []


def format_for_prompt(results: list[dict], min_score: float = 0.55) -> str:
    """Format file search hits for injection into the Jarvis system prompt."""
    hits = [r for r in results if r.get("score", 0) >= min_score]
    if not hits:
        return ""
    lines = ["[RELEVANT FILES — retrieved by semantic search]"]
    for r in hits:
        name  = r.get("file_name", "unknown")
        score = int(r.get("score", 0) * 100)
        text  = r.get("text", "").strip()[:350]
        lines.append(f"{name} ({score}% match): {text}")
    result = "\n".join(lines) + "\n\n"
    return result[:2000]


# ── Stats ─────────────────────────────────────────────────────────────────────

def get_stats() -> dict:
    """Return current index statistics for the dashboard."""
    try:
        with _lock:
            index = _load_index()
        files: dict[str, dict] = {}
        for entry in index:
            fp = entry.get("file_path", "")
            if fp not in files:
                files[fp] = {
                    "name":    entry.get("file_name", Path(fp).name),
                    "ext":     entry.get("file_ext", ""),
                    "chunks":  0,
                    "date":    entry.get("date", ""),
                    "size_kb": entry.get("size_kb", 0),
                }
            files[fp]["chunks"] += 1

        file_list = sorted(files.values(), key=lambda x: x["date"], reverse=True)
        return {
            "total_chunks":   len(index),
            "total_files":    len(files),
            "files":          file_list[:100],
            "watched_folders": [str(f) for f in DEFAULT_FOLDERS],
        }
    except Exception:
        return {"total_chunks": 0, "total_files": 0, "files": [], "watched_folders": []}


def index_custom_folder(folder_path: str, api_key: str) -> int:
    """Index any arbitrary folder path on the PC (validated, no system dirs)."""
    p = _safe_folder_path(folder_path)
    if not p.exists() or not p.is_dir():
        raise ValueError(f"Folder not found: {folder_path}")
    return scan_and_index([p], api_key)


def remove_file(file_path: str) -> bool:
    """Remove all chunks for a specific file from the index."""
    try:
        with _lock:
            index = _load_index()
            vecs  = _load_vectors()
            keep  = [i for i, e in enumerate(index) if e.get("file_path") != file_path]
            new_index = [index[i] for i in keep]
            new_vecs  = vecs[keep] if keep else np.empty((0, EMBED_DIM), dtype=np.float32)
            _save_index(new_index)
            _save_vectors(new_vecs)
        logger.info("Removed: %s", Path(file_path).name)
        return True
    except Exception as e:
        logger.error("remove_file error: %s", e)
        return False
